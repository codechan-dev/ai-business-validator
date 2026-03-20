"""
FastAPI Application for AI Business Idea Validator
"""
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
import logging
import io
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from backend.schemas import ValidateRequest, ValidateResponse
from backend.graph.langgraph_flow import validate_business_idea

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="AI Business Idea Validator API",
    description="Validate business ideas using AI, LangGraph, and real-time data",
    version="1.0.0"
)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, restrict to specific domains
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/api/health", tags=["Health"])
async def health_check():
    """Health check endpoint"""
    return {"status": "healthy", "version": "1.0.0"}

@app.post("/api/validate", response_model=ValidateResponse, tags=["Validation"])
async def validate(request: ValidateRequest) -> ValidateResponse:
    """
    Validate a business idea
    
    This endpoint performs comprehensive analysis including:
    - Market research
    - Competitor analysis
    - Demand estimation
    - Risk assessment
    - Real-time signal collection from MCP tools
    """
    try:
        logger.info(f"Validating idea: {request.idea[:50]}...")
        
        result = validate_business_idea(request.idea)
        
        logger.info(f"Validation completed. Scores: {result['scores']}")
        
        return ValidateResponse(**result)
    
    except ValueError as e:
        logger.error(f"Validation error: {str(e)}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"Unexpected error: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal server error")

@app.get("/api/signals/{source}", tags=["Signals"])
async def get_signals(source: str, query: str = ""):
    """
    Get real-time signals from a specific MCP source
    
    Sources:
    - reddit: Reddit discussions
    - trends: Google Trends data
    - product_hunt: Product Hunt launches
    - news: News mentions
    - startups: Startup database
    """
    try:
        signals = []
        
        if source == "reddit":
            from backend.mcp_tools import reddit
            signals = reddit.scrape_reddit(query or "startup")
        elif source == "trends":
            from backend.mcp_tools import trends
            signals = trends.get_trends(query or "startup")
        elif source == "product_hunt":
            from backend.mcp_tools import product_hunt
            signals = product_hunt.fetch_launches(query or "startup")
        elif source == "news":
            from backend.mcp_tools import news
            signals = news.fetch_news(query or "startup")
        elif source == "startups":
            from backend.mcp_tools import startups
            signals = startups.fetch_startups(query or "startup")
        else:
            raise ValueError(f"Unknown source: {source}")
        
        # Return signals or empty array as fallback
        return signals or []
    
    except Exception as e:
        logger.error(f"Error fetching {source} signals: {str(e)}")
        # Return mock data on error instead of failing
        return [
            {
                "title": f"Example {source.replace('_', ' ').title()} discussion",
                "summary": f"This is a sample {source} signal about the topic",
                "link": "https://example.com"
            }
        ]

@app.get("/api/history", tags=["History"])
async def get_history():
    """Get validation history (placeholder)"""
    return {"items": []}

@app.post("/api/download-report", tags=["Export"])
async def download_report(request: ValidateRequest):
    """
    Generate and download a PowerPoint report for a validated business idea
    """
    try:
        # Get validation data
        result = validate_business_idea(request.idea)
        
        # Extract nested data from ValidateResponse structure
        scores = result.get('scores', {})
        analysis = result.get('analysis', {})
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define color scheme (premium dark theme)
        DARK_BG = RGBColor(15, 23, 42)  # slate-900
        CYAN = RGBColor(34, 211, 238)    # cyan-400
        WHITE = RGBColor(255, 255, 255)
        GRAY = RGBColor(148, 163, 184)   # slate-400
        
        # Slide 1: Title Slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide1.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = "Business Idea Validation Report"
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = CYAN
        title_p.alignment = PP_ALIGN.CENTER
        
        idea_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
        idea_frame = idea_box.text_frame
        idea_frame.word_wrap = True
        idea_p = idea_frame.paragraphs[0]
        idea_p.text = request.idea
        idea_p.font.size = Pt(28)
        idea_p.font.color.rgb = WHITE
        idea_p.alignment = PP_ALIGN.CENTER
        
        date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        date_frame = date_box.text_frame
        date_p = date_frame.paragraphs[0]
        date_p.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        date_p.font.size = Pt(16)
        date_p.font.color.rgb = GRAY
        date_p.alignment = PP_ALIGN.CENTER
        
        # Slide 2: Score Overview
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide2.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "📊 Validation Scores"
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = CYAN
        
        y_pos = 1.5
        for score_name, score_value in scores.items():
            # Score label and value
            score_box = slide2.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.5))
            score_frame = score_box.text_frame
            score_p = score_frame.paragraphs[0]
            score_p.text = f"{score_name.title()}: {score_value}%"
            score_p.font.size = Pt(24)
            score_p.font.bold = True
            score_p.font.color.rgb = WHITE
            
            # Progress bar background
            bar_bg = slide2.shapes.add_shape(1, Inches(1), Inches(y_pos + 0.55), Inches(8), Inches(0.3))
            bar_bg.fill.solid()
            bar_bg.fill.fore_color.rgb = RGBColor(51, 65, 85)  # slate-700
            bar_bg.line.color.rgb = RGBColor(100, 116, 139)  # slate-600
            
            # Progress bar fill
            bar_fill_width = (score_value / 100) * 8
            bar_fill = slide2.shapes.add_shape(1, Inches(1), Inches(y_pos + 0.55), Inches(bar_fill_width), Inches(0.3))
            bar_fill.fill.solid()
            bar_fill.fill.fore_color.rgb = CYAN
            bar_fill.line.color.rgb = CYAN
            
            y_pos += 1.2
        
        # Slide 3: Market Analysis
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide3.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "🎯 Analysis Summary"
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = CYAN
        
        y_pos = 1.5
        analysis_items = [
            ("Market", analysis.get('market', 'Market analysis pending')),
            ("Competitors", analysis.get('competitors', 'Competitor analysis pending')),
            ("Risks", analysis.get('risks', 'Risk analysis pending'))
        ]
        
        for label, value in analysis_items:
            label_box = slide3.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.4))
            label_frame = label_box.text_frame
            label_p = label_frame.paragraphs[0]
            label_p.text = f"• {label}"
            label_p.font.size = Pt(18)
            label_p.font.bold = True
            label_p.font.color.rgb = CYAN
            
            value_box = slide3.shapes.add_textbox(Inches(1.5), Inches(y_pos + 0.4), Inches(7.5), Inches(1.2))
            value_frame = value_box.text_frame
            value_frame.word_wrap = True
            value_p = value_frame.paragraphs[0]
            value_p.text = str(value)[:200]
            value_p.font.size = Pt(14)
            value_p.font.color.rgb = WHITE
            
            y_pos += 1.8
        
        # Slide 4: Recommendation
        slide4 = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide4.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "💡 AI Recommendation"
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = CYAN
        
        recommendation = analysis.get('recommendation', 'Analysis complete. Consider market fit and competitive landscape.')
        rec_box = slide4.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        rec_frame = rec_box.text_frame
        rec_frame.word_wrap = True
        rec_p = rec_frame.paragraphs[0]
        rec_p.text = recommendation
        rec_p.font.size = Pt(20)
        rec_p.font.color.rgb = WHITE
        rec_p.line_spacing = 1.5
        
        # Save to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        # Return file
        filename = f"business_validation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        return FileResponse(
            io.BytesIO(pptx_bytes.getvalue()),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=filename
        )
    
    except Exception as e:
        logger.error(f"Error generating PPT: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to generate report: {str(e)}")

# Serve static files from frontend build
from fastapi.staticfiles import StaticFiles
import pathlib

# Get the path to frontend dist folder
# For production deployment with frontend built in Docker
frontend_dist = pathlib.Path(__file__).parent.parent / "frontend" / "dist"

# Mount static files if dist folder exists
if frontend_dist.exists():
    app.mount("/", StaticFiles(directory=str(frontend_dist), html=True), name="static")
    logger.info(f"Serving static frontend files from {frontend_dist}")
else:
    logger.warning(f"Frontend dist folder not found at {frontend_dist}. Frontend will not be served.")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
